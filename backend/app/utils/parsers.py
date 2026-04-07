"""
Raw text extraction from PDF, DOCX, and PPTX files.
Returns list of (page_number, text) tuples.
"""

import fitz  # PyMuPDF
import docx
from pptx import Presentation


def parse_pdf(file_bytes: bytes) -> list[tuple[int, str]]:
    """Extract text per page from a PDF."""
    pages = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text")
        if text.strip():
            pages.append((page_num, text))
    doc.close()
    return pages


def parse_docx(file_bytes: bytes) -> list[tuple[int, str]]:
    """Extract text from DOCX — treated as a single 'page'."""
    import io
    doc = docx.Document(io.BytesIO(file_bytes))
    full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [(1, full_text)]


def parse_pptx(file_bytes: bytes) -> list[tuple[int, str]]:
    """Extract text per slide from PPTX."""
    import io
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        combined = "\n".join(t for t in texts if t.strip())
        if combined:
            slides.append((slide_num, combined))
    return slides


SUPPORTED_PARSERS = {
    "application/pdf": parse_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": parse_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": parse_pptx,
}

SUPPORTED_EXTENSIONS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
}
